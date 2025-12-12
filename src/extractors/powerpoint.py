"""
PowerPoint extractor - Extract slide text, notes, and images from .pptx files
"""

import logging
from pathlib import Path
from typing import Optional

from extractors.base import BaseExtractor, ExtractionResult, ExtractionInterrupted
from utils.office_converter import OfficeConverter

logger = logging.getLogger(__name__)


class PowerPointExtractor(BaseExtractor):
    """Extract text and images from PowerPoint presentations (.pptx)"""
    
    SUPPORTED_EXTENSIONS = ['.pptx']
    
    def __init__(self, output_base_path: Path, extract_images: bool = False):
        super().__init__(output_base_path)
        self.pptx_available = self._check_pptx()
        self.extract_images = extract_images
        self.converter = OfficeConverter()
    
    def _check_pptx(self) -> bool:
        """Check if python-pptx is available"""
        try:
            import pptx
            return True
        except ImportError:
            logger.warning("python-pptx not available - PowerPoint extraction disabled")
            return False
    
    def can_extract(self, filepath: Path) -> bool:
        """Check if this extractor can handle the given file"""
        return filepath.suffix.lower() == '.pptx' and self.pptx_available
    
    def extract(self, filepath: Path, output_dir: Path) -> ExtractionResult:
        """
        Extract text and images from PowerPoint presentation
        
        Args:
            filepath: Path to PowerPoint file
            output_dir: Directory for extracted files
            
        Returns:
            ExtractionResult with extraction details
        """
        result = ExtractionResult(filepath)
        
        if not self.pptx_available:
            result.add_error("python-pptx not installed - cannot process PowerPoint files")
            return result
        
        try:
            import pptx
            
            logger.info(f"Extracting PowerPoint file: {filepath.name}")
            
            # Ensure output directory exists
            if not self.ensure_output_dir(output_dir):
                result.add_error(f"Failed to create output directory: {output_dir}")
                return result
            
            # Create subdirectory for images
            file_safe_name = self.sanitize_filename(filepath.name)
            
            # Open presentation
            prs = pptx.Presentation(filepath)
            
            result.metadata['slide_count'] = len(prs.slides)
            total_slides = len(prs.slides)
            logger.info(f"Presentation has {total_slides} slides")
            
            # Report initial substep
            self.report_substep(f"Extracting text from {total_slides} slides")
            
            # Extract text from all slides
            text_output = output_dir / f"{file_safe_name}.txt"
            text_content = self._extract_text(prs, result, total_slides)
            
            if text_content.strip():
                with open(text_output, 'w', encoding='utf-8') as f:
                    f.write(text_content)
                result.add_file(text_output)
                logger.info(f"Extracted text to {text_output.name}")
            else:
                result.add_warning("No text content found in presentation")
            
            # Check for interrupt before image extraction
            self.check_interrupted()
            
            # Extract slide images (snapshots)
            if self.extract_images:
                # Use LibreOffice converter for full slide screenshots
                images_dir = output_dir / f"{file_safe_name}_slides"
                
                # Check if LibreOffice is available
                if self.converter.soffice_path:
                    self.report_substep(f"Converting {total_slides} slides to images")
                    logger.info("Using LibreOffice for slide image extraction")
                    generated_files = self.converter.convert_to_png(filepath, images_dir)
                    
                    if generated_files:
                        image_count = len(generated_files)
                        for img in generated_files:
                            result.add_file(img)
                        result.metadata['slide_images_extracted'] = image_count
                        logger.info(f"Extracted {image_count} slide images via LibreOffice")
                    else:
                        logger.warning("LibreOffice conversion failed or produced no output")
                        result.add_warning("Slide image extraction failed")
                else:
                    logger.warning("LibreOffice not found - skipping slide image extraction")
                    result.add_warning("LibreOffice not found - cannot generate slide screenshots")
            else:
                 # We always want slide screenshots now, so if disabled by config, we log it.
                 # But previously I removed the 'else' block which kept "Image extraction disabled" message.
                 # Re-adding minimal logging if needed, or simply pass.
                 logger.info("Image extraction disabled for PowerPoint")
            if len(result.extracted_files) > 0:
                result.success = True
                logger.info(f"Successfully extracted data from {filepath.name}")
            else:
                result.add_warning("No data extracted from PowerPoint")
        
        except ExtractionInterrupted:
            # Re-raise interrupt exceptions so they propagate to the manager
            raise
            
        except Exception as e:
            error_msg = f"Failed to extract {filepath.name}: {str(e)}"
            logger.error(error_msg)
            result.add_error(error_msg)
        
        return result
    
    def _extract_text(self, prs, result: ExtractionResult, total_slides: int) -> str:
        """Extract all text from PowerPoint presentation"""
        try:
            text_parts = []
            
            text_parts.append(f"PowerPoint Presentation\n")
            text_parts.append(f"Total Slides: {len(prs.slides)}\n")
            text_parts.append(f"{'='*80}\n\n")
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                # Check for interrupt before each slide
                self.check_interrupted()
                
                # Report substep progress
                self.report_substep(f"Processing slide {slide_idx} of {total_slides}")
                
                text_parts.append(f"\n{'='*80}\n")
                text_parts.append(f"SLIDE {slide_idx}\n")
                text_parts.append(f"{'='*80}\n\n")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(f"{shape.text}\n")
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if hasattr(notes_slide, 'notes_text_frame'):
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            text_parts.append(f"\n--- NOTES ---\n")
                            text_parts.append(f"{notes_text}\n")
                
                text_parts.append("\n")
            
            return ''.join(text_parts)
        
        except ExtractionInterrupted:
            raise
            
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            result.add_warning(f"Text extraction error: {e}")
            return ""
